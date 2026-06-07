import io
import json
import logging
import posixpath
import re
import time
from datetime import datetime
from email.utils import parsedate_to_datetime
from subprocess import TimeoutExpired
from typing import Any
from urllib.parse import (
    parse_qsl,
    quote,
    unquote,
    urlencode,
    urlsplit,
    urlunsplit,
)

import requests
from bs4 import BeautifulSoup
from django.conf import settings
from django.db import connection, transaction
from django.db.models import Q
from pptx import Presentation
from PyPDF2 import PdfReader
from playwright.sync_api import TimeoutError as PlaywrightTimeoutError, sync_playwright
from tldextract import tldextract

from dbdb.core.models import (
    Acquisition, CitationUrl, CitationUrlContent, Organization,
    RepositoryInfo, System, SystemFeature, SystemVersion,
)
from dbdb.core.utils import spam
from dbdb.core.utils.git import get_git_commit_metadata
from dbdb.core.utils.spam import UnexpectedResponseError

LOG = logging.getLogger(__name__)

# --- Configuration ---

MAX_DOWNLOAD_BYTES = 60 * 1024 * 1024
REQUEST_TIMEOUT = 15 # seconds

SKIP_DOMAINS = {
    "//www.crunchbase.com/", # Recaptcha blocks
    "//twitter.com",
    "//www.bloomberg.com/",
    "//dl.acm.org/",
    "//dbdb.io/",
    "//www.linkedin.com/",
    "//docs.4d.com/",
    "//doc.4d.com/",
    '//git-wip-us.apache.org',
    '//angel.co/',
    '//www.cnet.com/', # HTML never renders?
}

SPAM_IGNORE_DOMAINS = {
    # "apache.org",
    "//en.wikipedia.org/",
    "//www.postgresql.org/",
    "//www.slideshare.net/", # LLM spam checker can't handle it?
    "//news.ycombinator.com/",
    "//books.google.com/",
    "//github.com/"
}

IGNORE_TITLES = map(str.lower, [
    "Redirecting to Google Groups",
    "PowerPoint Presentation",
    "PowerPoint 演示文稿",
    "Just a moment...",
    "File not found · GitHub",
    "Redirecting…",
])

REMOVE_TEXT = [
    "There was an error while loading. Please reload this page", # github
    " Uh oh! \n",  # github
]

# --- Exceptions ---

class UnsupportedContentTypeError(RuntimeError):
    pass


# --- Helpers ---

def _get_fragment(url: str) -> str:
    parts = urlsplit(url)
    return parts.fragment

def _extract_pdf_metadata(url:str, data: bytes, system: System | None = None) -> tuple[str | None, datetime | None]:
    """
    Extract title and oldest date (CreationDate or ModDate) from PDF metadata.

    Returns:
        tuple: (title, oldest_date) where title is a string or None,
               and oldest_date is a timezone-aware datetime or None
    """
    reader = PdfReader(io.BytesIO(data))
    meta = reader.metadata
    LOG.debug(f"meta: {meta}")

    # Extract title
    title = None
    if meta and meta.title:
        title = meta.title.strip() or None

    # Extract dates
    creation_date = None
    mod_date = None
    if meta:
        # Bad meta-data throws errors. We can just ignore it
        try:
            creation_date = meta.creation_date
        except:
            pass
        try:
            mod_date = meta.modification_date
        except:
            pass

    # Find the oldest date
    oldest_date = None
    if creation_date and mod_date:
        oldest_date = min(creation_date, mod_date)
    elif creation_date:
        oldest_date = creation_date
    elif mod_date:
        oldest_date = mod_date

    return title, oldest_date

def _extract_ppt_title(url:str, data: bytes, system: System | None = None) -> str | None:
    prs = Presentation(io.BytesIO(data))
    core = prs.core_properties
    title = None
    if core.title:
        title = core.title.strip()
    return title


def _parse_cache_control(header: str | None) -> dict[str, str | bool]:
    """
    Parse Cache-Control header into a structured dict.
    """
    if not header:
        return {}

    directives = {}
    for part in header.split(","):
        part = part.strip()
        if "=" in part:
            k, v = part.split("=", 1)
            directives[k.lower()] = v.strip('"')
        else:
            directives[part.lower()] = True
    return directives

def _extract_html_title(
        url: str,
        data: bytes,
        encoding: str | None = None,
        system: System | None = None,
        skip_spamcheck: bool = False,
        request_timeout: int | None = None,
    ) -> tuple[str | None, CitationUrl.Status | None]:

    if any(d in url for d in SPAM_IGNORE_DOMAINS):
        skip_spamcheck = True

    title = None
    html = _get_html_page(url, request_timeout=request_timeout)
    soup = BeautifulSoup(html, "html.parser")

    if html and (
        re.search(r'Visit\s+(?:<a[^>]*>)?cloudflare\.com(?:</a>)?\s+for more information', html, re.IGNORECASE)
        or "challenges.cloudflare.com" in html
    ):
        LOG.debug(f"Cloudflare challenge page detected for {url}")
        return None, CitationUrl.Status.IGNORE

    if not skip_spamcheck:
        # <noscript> content is hidden by browsers when JS is enabled; strip it
        # so the LLM doesn't see "You need to enable JavaScript to run this app."
        for tag in soup.find_all("noscript"):
            tag.decompose()

        # Extract the text from the HTML and clean up the newlines and spaces
        # This is wasted space in our prompt context
        text_words = soup.get_text(" ")
        text_words = re.sub(r"(?:\t){1,}", " ", text_words, flags=re.DOTALL)
        text_words = re.sub(r"(?: ){2,}", " ", text_words, flags=re.DOTALL)
        text_words = re.sub(r"(?: \n){2,}", " \n", text_words, flags=re.DOTALL)
        text_words = re.sub(r"(?:\n \n)", " \n", text_words, flags=re.DOTALL)

        for s in REMOVE_TEXT:
            text_words = text_words.replace(s, "")

        # Make sure there is at least something for us to look
        # at when we use the spam checker
        if len(text_words) > 0:
            attempts = 3
            temperature = 0.0
            model = settings.CRAWLER_SPAM_CHECKER_MODEL
            is_spam = None
            while attempts > 0:
                attempts -= 1
                try:
                    is_spam = spam.is_spam(text_words, system,
                                           timeout=request_timeout,
                                           temperature=temperature,
                                           model=model
                    )
                    break
                except UnexpectedResponseError as e:
                    # Ignore if we get a weird LLM response
                    LOG.error(e)
                    if attempts == 0: raise e
                    pass
                model = settings.CRAWLER_SPAM_CHECKER_FALLBACK_MODEL_A if attempts % 2 == 0 else settings.CRAWLER_SPAM_CHECKER_FALLBACK_MODEL_B
                temperature = max(temperature - 0.1, 0.0)
            if is_spam is not None and is_spam:
                return None, CitationUrl.Status.SPAM

    if soup.title and soup.title.string:
        title = soup.title.string.strip()

    return title, None

def _get_html_page(
        url,
        render_wait: float = 10,
        request_timeout: int | None = None) -> str | None:
    timeout_ms = int((request_timeout or render_wait) * 1000)

    html_source = None
    with sync_playwright() as pw:
        browser = pw.firefox.launch(headless=True)
        try:
            context = browser.new_context(
                user_agent=settings.CRAWLER_USER_AGENT,
                # Mask automation signals: JS property and HTTP header
                java_script_enabled=True,
            )
            # Hide navigator.webdriver before any page script runs
            context.add_init_script(
                "Object.defineProperty(navigator, 'webdriver', {get: () => undefined})"
            )
            page = context.new_page()
            try:
                # wait_until='networkidle' blocks until there are no more than
                # 0 in-flight network requests for 500 ms — ideal for SPAs.
                page.goto(url, wait_until="networkidle", timeout=timeout_ms)
                LOG.debug(f"Page loaded (networkidle): {page.title()!r}")
                html_source = page.content()
            except PlaywrightTimeoutError:
                LOG.debug(f"Timed out waiting for networkidle on {url}; capturing partial page source")
                try:
                    html_source = page.content()
                except Exception:
                    pass
            finally:
                page.close()
        finally:
            browser.close()

    return html_source

# --- Main API ---

def _extract_wikipedia_metadata(
    data: bytes,
    encoding: str | None = None,
) -> tuple[str | None, datetime | None]:
    """Extract title and last_modified from a Wikipedia page without JS rendering.

    Title comes from <meta property="og:title">; last_modified comes from the
    dateModified field in the embedded application/ld+json script block.
    """
    html = data.decode(encoding or 'utf-8', errors='replace')
    soup = BeautifulSoup(html, 'html.parser')

    title = None
    og_title = soup.find('meta', attrs={'property': 'og:title'})
    if og_title and og_title.get('content'):
        title = og_title['content'].strip() or None

    last_modified = None
    for script in soup.find_all('script', attrs={'type': 'application/ld+json'}):
        try:
            ld = json.loads(script.string or '')
            date_str = ld.get('dateModified')
            if date_str:
                last_modified = datetime.fromisoformat(date_str)
                break
        except (json.JSONDecodeError, ValueError, AttributeError):
            pass

    return title, last_modified


TEXT_CONTENT_MAX_CHARS = 8_000


def fetch_url_metadata(
    url: str,
    *,
    system: System | None = None,
    citation_url: CitationUrl | None = None,
    skip_spamcheck: bool = False,
    request_timeout: int | None = None,
    if_none_match: str | None = None,
    if_modified_since: datetime | None = None,
    allow_redirects: bool = False,
    redirect_ctr: int = 0,
) -> dict[str, Any]:

    headers = {"User-Agent": settings.CRAWLER_USER_AGENT}
    if if_none_match:
        headers["If-None-Match"] = if_none_match
    if if_modified_since:
        headers["If-Modified-Since"] = if_modified_since.strftime(
            "%a, %d %b %Y %H:%M:%S GMT"
        )

    if not request_timeout:
        request_timeout = REQUEST_TIMEOUT

    # Special Case: Busted Domains
    if any(s in url for s in SKIP_DOMAINS):
        LOG.debug(f"Skipping '{url}' because it is from a domain to ignore")
        return {
            "url": url,
            "status-code": None,
            "content-type": None,
            "content-length": None,
            "status": CitationUrl.Status.IGNORE,
            "title": None,
            "etag": None,
            "last-modified": None,
            "cache-control": None,
            "revalidate": None
        }

    # Special Case: Github Commit
    m = re.match(
        r"https://github\.com/(?P<owner>[^/]+)/(?P<repo>[^/.]+)/commit/(?P<commit>[^/.]+)/?$",
        url,
    )
    if m and "#diff-" not in url:
        owner = m.group("owner")
        repo = m.group("repo")
        commit_id = m.group("commit")
        repo_url = f"https://github.com/{owner}/{repo}.git"
        try:
            # If the Github repo is no longer available, then this will time out
            # So we'll just treat it like a regular page and record the 404 status-code
            title, last_modified = get_git_commit_metadata(repo_url, commit_id, timeout=request_timeout)
            return {
                "url": url,
                "status-code": 200,
                "content-type": None,
                "content-length": None,
                "status": CitationUrl.Status.VALID,
                "title": title,
                "etag": None,
                "last-modified": last_modified,
                "cache-control": None,
                "revalidate": None
            }
        except TimeoutExpired as e:
            LOG.debug(f"Timeout: {e}")
            pass

    title = None
    status: CitationUrl.Status = CitationUrl.Status.UNKNOWN
    raw_content = ''
    clean_text = ''
    _html_encoding = 'utf-8'

    LOG.debug(f"Fetching '{url}'\n -> Headers: {headers}" )
    with requests.get(
        url,
        stream=True,
        timeout=REQUEST_TIMEOUT,
        headers=headers,
        allow_redirects=True if redirect_ctr > 8 else allow_redirects,
    ) as resp:
        status_code = resp.status_code
        LOG.debug(f"{url}\nstatus_code={status_code}")

        # If we get redirected, then recursively call ourselves
        # with allowing the redirect so that we can get the new URL
        if status_code in (301, 302, 307, 308):
            new_url = resp.headers['Location']
            if not new_url.startswith("http"):
                orig_extracted = tldextract.extract(url)
                new_url = "https://" + orig_extracted.fqdn + new_url
            # Don't normalize because we may get stuck in an infinite loop
            # if the server-side adds back the trailing slash
            normalized_url = normalize_url(new_url)
            if normalized_url + '/' != new_url: new_url = normalized_url

            # Add back the fragment after normalize the redirect
            fragment = _get_fragment(url)
            if fragment: new_url += f"#{fragment}"

            LOG.debug(f"Redirect: {new_url}")
            result = fetch_url_metadata(
                new_url,
                system=system,
                skip_spamcheck=skip_spamcheck,
                request_timeout=request_timeout,
                if_none_match=if_none_match,
                if_modified_since=if_modified_since,
                allow_redirects=False,
                redirect_ctr=redirect_ctr+1,
            )
            if not result: return None
            if result["status"] != CitationUrl.Status.SPAM:
                # Only update the URL redirect if the domains are the same.
                # Otherwise they will redirect to a spam cite and we lose the original URL
                orig_extracted = tldextract.extract(url)
                new_extracted = tldextract.extract(new_url)
                if orig_extracted.domain == new_extracted.domain and \
                    orig_extracted.suffix == new_extracted.suffix:
                    LOG.debug(f"Updating URL: status={result['status']} / {new_url}")
                    result["url"] = new_url
            else:
                result["url"] = url
            return result

        content_type = (
            resp.headers.get("Content-Type", "").split(";")[0].lower() or None
        )
        content_length = resp.headers.get("Content-Length", None)
        etag = resp.headers.get("ETag")
        _html_encoding = resp.encoding or 'utf-8'

        last_modified_hdr = resp.headers.get("Last-Modified")
        try:
            last_modified = (
                parsedate_to_datetime(last_modified_hdr)
                if last_modified_hdr
                else None
            )
        except (TypeError, ValueError):
            last_modified = None

        cache_control = _parse_cache_control(
            resp.headers.get("Cache-Control")
        )

        # --- Short-circuit on 304 ---
        if status_code == 304:
            return {
                "url": url,
                "status": CitationUrl.Status.VALID,
                "status-code": status_code,
                "content-type": content_type,
                "content-length": content_length,
                "title": None,
                "etag": etag,
                "last-modified": last_modified,
                "cache-control": cache_control,
                "revalidate": {
                    "if-none-match": etag or if_none_match,
                    "if-modified-since": last_modified or if_modified_since,
                },
            }

        # --- Download body ---
        data = bytearray()
        for chunk in resp.iter_content(chunk_size=8192):
            data.extend(chunk)
            if len(data) > MAX_DOWNLOAD_BYTES:
                raise RuntimeError(f"Download exceeds size limit [#bytes={len(data)}]")
        data = bytes(data)
        if status_code != 404 and len(data) == 0:
            LOG.error(f"Unexpected empty contents for '{url}'")
            status = CitationUrl.Status.IGNORE

    if status == CitationUrl.Status.UNKNOWN:
        status = CitationUrl.Status.DEAD
        if (200 <= status_code < 300) or status_code == 307:
            status = CitationUrl.Status.VALID

    # --- Content-Type dispatch ---

    if status_code != 404 and status != CitationUrl.Status.IGNORE:
        # PDF
        if content_type == "application/pdf":
            title, pdf_last_modified = _extract_pdf_metadata(url, data)
            if pdf_last_modified: last_modified = pdf_last_modified
            try:
                reader = PdfReader(io.BytesIO(data))
                raw_content = '\n'.join(page.extract_text() or '' for page in reader.pages)
                clean_text = raw_content[:TEXT_CONTENT_MAX_CHARS]
            except Exception:
                pass

        # PPTX
        elif content_type in {
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }:
            title = _extract_ppt_title(url, data)

        # HTML — Wikipedia: parse raw bytes directly, skip JS rendering
        elif content_type in {"text/html", "application/xhtml+xml"} and "wikipedia.org" in url:
            wiki_title, wiki_last_modified = _extract_wikipedia_metadata(data, _html_encoding)
            if wiki_title:
                title = wiki_title
            if wiki_last_modified:
                last_modified = wiki_last_modified
            raw_content = data.decode(_html_encoding, errors='replace')
            try:
                soup = BeautifulSoup(data, 'html.parser')
                clean_text = re.sub(r'\s+', ' ', soup.get_text(' ')).strip()[:TEXT_CONTENT_MAX_CHARS]
            except Exception:
                pass

        # HTML
        elif content_type in {"text/html", "application/xhtml+xml"}:
            title, page_status = _extract_html_title(url, data,
                                                     encoding=_html_encoding,
                                                     skip_spamcheck=skip_spamcheck,
                                                     system=system,
                                                     request_timeout=request_timeout)
            if page_status is not None:
                status = page_status
            raw_content = data.decode(_html_encoding, errors='replace')
            try:
                soup = BeautifulSoup(data, 'html.parser')
                for tag in soup.find_all(['script', 'style', 'noscript']):
                    tag.decompose()
                clean_text = re.sub(r'\s+', ' ', soup.get_text(' ')).strip()[:TEXT_CONTENT_MAX_CHARS]
            except Exception:
                pass

    # Minor cleaning...
    if title is not None and title:
        title = title.replace("\n", " ")
        title = re.sub(r' {2,}', ' ', title)
        if title and title.strip() == url: title = None
        if title and title.startswith("https://"): title = None
        if title and title.lower() in IGNORE_TITLES: title = None
        if title and status_code in [403, 404]: title = None

    if citation_url is not None and (raw_content or clean_text):
        CitationUrlContent.objects.update_or_create(
            citation=citation_url,
            defaults={'raw': raw_content, 'text': clean_text},
        )

    return {
        "url": url,
        "status": status,
        "status-code": status_code,
        "content-type": content_type,
        "content-length": content_length,
        "title": title,
        "etag": etag,
        "last-modified": last_modified,
        "cache-control": cache_control,
        "revalidate": {
            "if-none-match": etag,
            "if-modified-since": last_modified,
        },
        "raw": raw_content,
        "text": clean_text,
    }

def merge_citations(merge_to: CitationUrl, merge_from: list[CitationUrl]) -> CitationUrl:
    """
    Merge all references to each CitationUrl in *merge_from* into *merge_to*,
    then delete the now-unreferenced merge_from rows.

    Handles:
      - M2M through-tables (SystemFeature.citations, SystemVersion citation M2Ms)
      - FK fields on Organization (url, linkedin_url)
      - FK field on Acquisition (citation)
      - FK fields on SystemVersion (system_url, docs_url, sourcerepo_url, wikipedia_url)
      - OneToOne field on RepositoryInfo (sourcerepo_url) — if merge_to already has
        a RepositoryInfo, the duplicate is deleted; otherwise it is reassigned.
    """
    LOG.debug(f"Merging {len(merge_from)} citation(s) into {merge_to}")

    # M2M through-tables: (table, owner_column)
    m2m_tables = [
        ("core_systemfeature_citations",             "systemfeature_id"),
        ("core_systemversion_description_citations", "systemversion_id"),
        ("core_systemversion_end_year_citations",    "systemversion_id"),
        ("core_systemversion_history_citations",     "systemversion_id"),
        ("core_systemversion_start_year_citations",  "systemversion_id"),
        ("core_feature_citations",                   "feature_id"),
        ("core_featureoption_citations",             "featureoption_id"),
        ("core_attribute_citations",                 "attribute_id"),
        ("core_attributeoption_citations",           "attributeoption_id"),
        ("core_docpage_citations",                   "docpage_id"),
    ]

    from_ids = [c.id for c in merge_from]
    all_ids  = from_ids + [merge_to.id]

    with transaction.atomic():
        # ── M2M tables ────────────────────────────────────────────────────
        for table, owner_col in m2m_tables:
            with connection.cursor() as cursor:
                placeholders = ', '.join(['%s'] * len(all_ids))
                cursor.execute(
                    f"SELECT id, {owner_col}, citationurl_id "
                    f"FROM {table} WHERE citationurl_id IN ({placeholders})",
                    all_ids,
                )
                rows = cursor.fetchall()

                # Pairs (owner_id, citationurl_id) already pointing at merge_to
                keep_pairs: set[tuple[int, int]] = {
                    (row[1], row[2]) for row in rows if row[2] == merge_to.id
                }
                to_update: list[tuple[int, int]] = []
                to_delete: list[tuple[int, int]] = []

                for _row_id, owner_id, cit_id in rows:
                    if cit_id == merge_to.id:
                        continue
                    if (owner_id, merge_to.id) in keep_pairs:
                        to_delete.append((owner_id, cit_id))
                    else:
                        to_update.append((owner_id, cit_id))
                        keep_pairs.add((owner_id, merge_to.id))

                for owner_id, cit_id in to_update:
                    cursor.execute(
                        f"UPDATE {table} SET citationurl_id = %s "
                        f"WHERE citationurl_id = %s AND {owner_col} = %s",
                        [merge_to.id, cit_id, owner_id],
                    )
                    LOG.info("Updated %d row(s) in %s", cursor.rowcount, table)

                for owner_id, cit_id in to_delete:
                    cursor.execute(
                        f"DELETE FROM {table} "
                        f"WHERE citationurl_id = %s AND {owner_col} = %s",
                        [cit_id, owner_id],
                    )
                    LOG.info("Deleted %d row(s) in %s", cursor.rowcount, table)

        # ── FK fields ─────────────────────────────────────────────────────
        for citation in merge_from:
            Organization.objects.filter(url=citation).update(url=merge_to)
            Organization.objects.filter(linkedin_url=citation).update(linkedin_url=merge_to)

            Acquisition.objects.filter(citation=citation).update(citation=merge_to)

            SystemVersion.objects.filter(system_url=citation).update(system_url=merge_to)
            SystemVersion.objects.filter(docs_url=citation).update(docs_url=merge_to)
            SystemVersion.objects.filter(sourcerepo_url=citation).update(sourcerepo_url=merge_to)
            SystemVersion.objects.filter(wikipedia_url=citation).update(wikipedia_url=merge_to)

            # RepositoryInfo is OneToOne with CASCADE — reassign or drop duplicate
            try:
                ri = RepositoryInfo.objects.get(sourcerepo_url=citation)
                if RepositoryInfo.objects.filter(sourcerepo_url=merge_to).exists():
                    ri.delete()
                    LOG.info("Deleted duplicate RepositoryInfo for %s", citation.url)
                else:
                    ri.sourcerepo_url = merge_to
                    ri.save()
                    LOG.info("Reassigned RepositoryInfo from %s to %s", citation.url, merge_to.url)
            except RepositoryInfo.DoesNotExist:
                pass

        # ── Delete the now-unreferenced source citations ───────────────────
        deleted, _ = CitationUrl.objects.filter(id__in=from_ids).delete()
        LOG.info("Deleted %d CitationUrl(s): %s", deleted, from_ids)

    return merge_to

def get_systems(c: CitationUrl,
                current_only: bool | None = False) -> list[System]:
    """
    Return the list of systems that reference this CitationURL
    :param citation:
    :return:
    """

    # Check the SystemFeatures first
    # features = SystemFeature.objects.filter(citations__in=[c])
    features = SystemFeature.objects.filter(citations__in=[c])
    if current_only:
        features = features.filter(version__is_current=True)
    systems = set([sf.version.system for sf in features])

    # Then check all CitationUrl fields on SystemVersions
    versions = SystemVersion.objects.filter(
        # M2M citation fields
        Q(description_citations__in=[c]) |
        Q(start_year_citations__in=[c]) |
        Q(end_year_citations__in=[c]) |
        Q(history_citations__in=[c]) |
        # FK URL fields
        Q(system_url=c) |
        Q(docs_url=c) |
        Q(sourcerepo_url=c) |
        Q(wikipedia_url=c) |
        # Developer org URLs
        Q(developer_orgs__url=c) |
        Q(developer_orgs__linkedin_url=c) |
        # Acquisition citation
        Q(acquisitions__citation=c)
    ).distinct()
    if current_only:
        versions = versions.filter(is_current=True)
    systems.update(v.system for v in versions)

    return list(systems)



def normalize_url(url: str) -> str:
    """
    Normalize a URL for comparison/deduplication while preserving fragments.
    """

    # HACK: Remove any URLs that already link to archive.org
    m = re.match(r"https://web\.archive\.org/web/.*?/(?P<url>http[s]?:/[/]?.*?)$", url, re.IGNORECASE)
    if m:
        url = m.group('url').strip()
        LOG.debug(f"Removed archive.org prefix: {url}")

    parts = urlsplit(url)

    # 1. Lowercase scheme and hostname
    scheme = parts.scheme.lower()
    hostname = parts.hostname.lower() if parts.hostname else ""

    # 2. Remove default ports
    try:
        port = parts.port
    except ValueError:
        port = 443 if scheme == "https" else 80
    except:
        raise
    if port and not (
        (scheme == "http" and port == 80)
        or (scheme == "https" and port == 443)
    ):
        netloc = f"{hostname}:{port}"
    else:
        netloc = hostname

    # Preserve userinfo if present
    if parts.username:
        userinfo = parts.username
        if parts.password:
            userinfo += f":{parts.password}"
        netloc = f"{userinfo}@{netloc}"
    # HACK: We sometimes got two dots at the end of a netloc?
    if netloc.endswith(".."):
        netloc = netloc[:-2]

    # 3. Normalize path
    path = unquote(parts.path)
    path = posixpath.normpath(path)

    # If there is no path, then this will add a dot. Remove that if we only have a dot
    if path in {".", "/"}: path = ""
    # If there is a path and it doesn't start with '/', then add the slash
    # We will purposely exclude a path for root domain URLs
    if path and not path.startswith("/"):
        path = "/" + path
    if path != "/" and path.endswith("/"):
        path = path[:-1]

    # Re-encode path
    path = quote(path, safe="/~:@")

    # 4. Normalize query (sorted)
    query_params = parse_qsl(parts.query, keep_blank_values=True)
    query = urlencode(sorted(query_params), doseq=True)

    # 5. Preserve fragment
    fragment = parts.fragment

    LOG.debug(f"Normalize: {url}\n+ scheme: {scheme}\n+ netloc: {netloc}\n+ path: {path}\n+ query: {query}\n+ fragment: {fragment}")
    return urlunsplit((scheme, netloc, path, query, fragment))
